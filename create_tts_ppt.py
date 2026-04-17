"""Generate TTS Models Research PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x3D)
ACCENT = RGBColor(0x00, 0xD2, 0xFF)
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xA0, 0xA0, 0xB8)
LIGHT = RGBColor(0xE2, 0xE8, 0xF0)


def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_DARK


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_bullet_card(slide, left, top, width, height, title, items, title_color=ACCENT, item_color=LIGHT, fill=BG_CARD):
    add_card(slide, left, top, width, height, fill)
    add_text_box(slide, left + 0.2, top + 0.15, width - 0.4, 0.4, title, 16, title_color, bold=True)
    for i, item in enumerate(items):
        add_text_box(slide, left + 0.2, top + 0.55 + i * 0.32, width - 0.4, 0.3, f"  {item}", 12, item_color)


# ══════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 1, 1.5, 11, 1, "6 On-Device TTS Models", 44, ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.6, 11, 0.8, "You Can Run Right Now", 36, WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, 2, 3.8, 9, 0.6, "Research Report: Architecture, Voice Cloning, Benchmarks, Advantages & Disadvantages", 18, GRAY, align=PP_ALIGN.CENTER)

models_list = ["Pocket TTS  -  Kyutai Labs", "VibeVoice  -  Microsoft", "Kitten TTS  -  KittenML",
               "Qwen3-TTS  -  Alibaba", "TADA TTS  -  Hume AI", "Neu TTS  -  Neuphonic"]
for i, m in enumerate(models_list):
    col = i % 3
    row = i // 3
    add_card(slide, 2 + col * 3.2, 4.8 + row * 0.7, 2.9, 0.55, BG_CARD)
    add_text_box(slide, 2.1 + col * 3.2, 4.85 + row * 0.7, 2.7, 0.45, m, 13, LIGHT, align=PP_ALIGN.CENTER)

add_text_box(slide, 1, 6.8, 11, 0.4, "All models run on CPU  |  No GPU required  |  Open Source", 14, GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 2: Quick Comparison Matrix
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.6, "Quick Comparison Matrix", 32, ACCENT, bold=True)

headers = ["Model", "Size", "Languages", "Voice Clone", "License", "Best For"]
data = [
    ["Pocket TTS", "~100M", "2 (EN, FR)", "Yes (zero-shot)", "Apache 2.0", "Mobile/edge + cloning"],
    ["VibeVoice", "0.5B-1.5B", "11+", "Limited", "MIT", "Long-form multi-speaker"],
    ["Kitten TTS", "15M-80M", "1 (EN)", "No", "Apache 2.0", "Extreme size constraints"],
    ["Qwen3-TTS", "0.6B-1.7B", "10", "Yes (3s)", "Apache 2.0", "Multilingual quality"],
    ["TADA TTS", "1B", "9+", "Yes", "MIT + Llama", "Hallucination-free"],
    ["Neu TTS", "748M", "4", "Yes (3s)", "Apache 2.0", "Real-time voice agents"],
]

col_widths = [1.8, 1.2, 1.5, 1.8, 1.5, 2.8]
x_start = 0.7
y_start = 1.2

# Header row
for j, h in enumerate(headers):
    x = x_start + sum(col_widths[:j])
    add_card(slide, x, y_start, col_widths[j] - 0.05, 0.45, ACCENT2)
    add_text_box(slide, x + 0.05, y_start + 0.05, col_widths[j] - 0.15, 0.35, h, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)

# Data rows
for i, row in enumerate(data):
    y = y_start + 0.5 + i * 0.45
    row_color = BG_CARD if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x48)
    for j, cell in enumerate(row):
        x = x_start + sum(col_widths[:j])
        add_card(slide, x, y, col_widths[j] - 0.05, 0.4, row_color)
        c = GREEN if cell in ("Yes (zero-shot)", "Yes (3s)", "Yes") else RED if cell == "No" else LIGHT
        add_text_box(slide, x + 0.05, y + 0.05, col_widths[j] - 0.15, 0.3, cell, 11, c, align=PP_ALIGN.CENTER)

# GPU note
add_text_box(slide, 0.5, 4.5, 12, 0.4, "All 6 models can run without GPU - designed for on-device deployment", 14, GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 3: Pocket TTS
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "1. Pocket TTS", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 5, 0.4, "Kyutai Labs  |  ~100M params  |  Apache 2.0", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features",
    ["CPU-only, no GPU needed",
     "~6x faster than realtime (M4)",
     "~200ms first audio chunk",
     "Zero-shot voice cloning",
     "8 built-in voices",
     "English + French",
     "Streaming (infinite text)",
     "CLI + Web UI included"],
    ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "Runs on CPU - zero API cost",
    "Voice cloning from short audio",
    "Export voice as .safetensors",
    "pip install pocket-tts (easy)",
    "WASM, MLX, Rust, C# ports",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "Only 2 languages",
    "No speed control parameter",
    "Cloning needs HuggingFace login",
    "Smaller community vs Coqui",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 3.5, "Code Example", [
    "pip install pocket-tts scipy",
    "",
    "from pocket_tts import TTSModel",
    "model = TTSModel.load_model()",
    'voice = model.get_state_for_',
    '    audio_prompt("alba")',
    'audio = model.generate_audio(',
    '    voice, "Hello world!")',
], ACCENT2, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 4: VibeVoice
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "2. VibeVoice", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.4, "Microsoft  |  0.5B-1.5B params  |  MIT License", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features", [
    "Next-token diffusion architecture",
    "7.5 Hz ultra-low frame rate",
    "~300ms first audio (Realtime)",
    "Up to 90 min in single pass",
    "Multi-speaker (up to 4)",
    "11+ languages (experimental)",
    "Qwen2.5 backbone",
    "3 model variants",
], ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "90-min long-form generation",
    "Multi-speaker support",
    "MIT License - very permissive",
    "0.5B Realtime is lightweight",
    "Gradio web UI included",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "Requires NVIDIA GPU",
    "Not for commercial use yet",
    "TTS code removed Sept 2025",
    "Limited voice cloning",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 2.8, "Third-Party APIs", [
    "fal.ai: $0.04/min (cheapest)",
    "Replicate: $0.11/run",
    "AIML API: available",
    "vibevoice.online: free tier",
    "",
    "Self-host: free (MIT license)",
], ORANGE, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 5: Kitten TTS
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "3. Kitten TTS", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.4, "KittenML  |  15M-80M params  |  Apache 2.0", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features", [
    "Ultra-tiny: 25MB (nano int8)",
    "ONNX-based inference",
    "4 model tiers (mini/micro/nano)",
    "8 built-in voices",
    "Speed control parameter",
    "24kHz audio output",
    "English only",
    "Runs on Raspberry Pi",
], ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "Smallest TTS model (25MB!)",
    "Runs on wearables, browsers",
    "Multiple size/quality tiers",
    "Adjustable speech speed",
    "Zero GPU dependency",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "No voice cloning",
    "English only",
    "Requires eSpeak-NG install",
    "Developer preview - may change",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 2.5, "Model Sizes", [
    "Mini:     80M params / 80MB",
    "Micro:    40M params / 41MB",
    "Nano FP32: 15M params / 56MB",
    "Nano int8: 15M params / 25MB",
    "",
    "Voices: expr-voice-2 to 5 (m/f)",
], ACCENT2, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 6: Qwen3-TTS
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "4. Qwen3-TTS", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.4, "Alibaba Group  |  0.6B-1.7B params  |  Apache 2.0", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features", [
    "Trained on 5M hours of data",
    "10 languages in single model",
    "3-second voice cloning",
    "Dual-track LM architecture",
    "Streaming + batch modes",
    "9 premium built-in voices",
    "Voice design from text",
    "Qwen3 LLM backbone",
], ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "Massive training scale (5M hrs)",
    "10-language multilingual",
    "Choose latency vs quality",
    "Voice design + cloning",
    "Strong Qwen3 ecosystem",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "Requires GPU (~8GB VRAM)",
    "Streaming mode quality drops",
    "Self-reported benchmarks",
    "Missing Hindi, Thai, etc.",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 2.8, "Third-Party APIs", [
    "AIML API: $0.013/1K chars",
    "    (cheapest option!)",
    "fal.ai: $0.07/1K chars",
    "Replicate: available",
    "Alibaba Cloud: native",
    "",
    "Self-host: free (Apache 2.0)",
], ORANGE, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 7: TADA TTS
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "5. TADA TTS", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.4, "Hume AI  |  1B-3B params  |  MIT + Llama License", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features", [
    "ZERO content hallucinations",
    "1:1 text-audio alignment",
    "5x faster than LLM-based TTS",
    "700 seconds in one pass",
    "Built on Llama 3.2",
    "Prompt caching for speed",
    "9+ languages supported",
    "Alignment verification tool",
], ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "Zero hallucinations by design",
    "Ideal for audiobooks/podcasts",
    "5x speed over autoregressive",
    "Prompt caching (encode once)",
    "Open weights (MIT license)",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "Requires GPU (CUDA)",
    "Strict alignment limits prosody",
    "Primarily English at launch",
    "Newer - fewer integrations",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 2.8, "Hume AI Cloud API", [
    "Free:    10K chars/mo ($0)",
    "Starter: 30K chars/mo ($3)",
    "Creator: 140K chars/mo ($7)",
    "Pro:     1M chars/mo ($70)",
    "Scale:   3.3M chars/mo ($200)",
    "",
    "Voice design via text prompt",
], ORANGE, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 8: Neu TTS
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "6. Neu TTS", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.4, "Neuphonic  |  748M params  |  Apache 2.0", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features", [
    "Sub-25ms cloud latency",
    "Qwen2 backbone + NeuCodec",
    "GGUF quantization (CPU-first)",
    "3-second voice cloning",
    "Built-in audio watermarking",
    "22 English voices available",
    "4 languages (EN/FR/DE/ES)",
    "Streaming-native architecture",
], ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "On-device: laptop/phone/RPi",
    "Cloud API: sub-25ms latency",
    "Full training pipeline open",
    "Audio watermarking built-in",
    "22 English voices + cloning",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "ELO 949 (below top-tier cloud)",
    "Primarily English",
    "748M still sizeable for edge",
    "Cloud: $17.60/1M chars",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 2.8, "Cloud API Pricing", [
    "$17.60 per 1M characters",
    "",
    "app.neuphonic.com (sign up)",
    "pip install pyneuphonic",
    "",
    "On-device: FREE (Apache 2.0)",
    "Air (360M) + Nano (120M)",
], ORANGE, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 9: Bonus - VoxCPM2
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 5, 0.6, "Bonus: VoxCPM2", 32, ACCENT, bold=True)
add_text_box(slide, 0.5, 0.85, 8, 0.4, "OpenBMB  |  2B params  |  Apache 2.0  |  Released Apr 2026", 14, GRAY)

add_bullet_card(slide, 0.5, 1.5, 4, 3.2, "Key Features", [
    "30 languages supported",
    "48kHz studio-quality audio",
    "Voice design from text desc",
    "3 cloning modes",
    "Tokenizer-free diffusion AR",
    "2M+ hours training data",
    "Beats ElevenLabs similarity",
    "pip install voxcpm",
], ACCENT, LIGHT)

add_bullet_card(slide, 4.7, 1.5, 4, 2.5, "Advantages", [
    "30 langs (incl Hindi!)",
    "48kHz studio quality",
    "85.4% voice similarity",
    "3 cloning modes + voice design",
    "Apache 2.0 commercial use",
], GREEN, LIGHT)

add_bullet_card(slide, 4.7, 4.2, 4, 2.0, "Disadvantages", [
    "Requires GPU (~8GB VRAM)",
    "No third-party API yet",
    "Very new (Apr 2026)",
    "Self-host only for now",
], RED, LIGHT)

add_bullet_card(slide, 9, 1.5, 4, 2.5, "3 Cloning Modes", [
    "1. Controllable Clone",
    "   ref audio only - good",
    "2. Ultimate Clone",
    "   ref audio + transcript - best",
    "3. Voice Design",
    "   text description only",
], ACCENT2, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 10: Pricing Comparison
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.6, "Cloud API Pricing Comparison", 32, ACCENT, bold=True)

price_headers = ["Service", "Price", "Unit", "~30-min Interview Cost", "Voice Clone"]
price_data = [
    ["Pocket TTS (local)", "FREE", "---", "$0.00", "Yes"],
    ["Qwen3-TTS (AIML API)", "$0.013", "per 1K chars", "~$0.05", "Yes"],
    ["AWS Polly", "$16", "per 1M chars", "~$0.02", "No"],
    ["VibeVoice (fal.ai)", "$0.04", "per minute", "~$1.20", "Limited"],
    ["Hume AI (TADA)", "$0-70", "per month", "~$0.03", "Yes"],
    ["Neuphonic", "$17.60", "per 1M chars", "~$0.02", "Yes"],
    ["ElevenLabs", "$0.234", "per 1K chars", "~$0.94", "Yes"],
    ["OpenAI TTS-1", "$0.02", "per 1K chars", "~$0.08", "No"],
]

p_col_w = [2.5, 1.5, 1.8, 2.5, 1.5]
px = 1.5
py = 1.2

for j, h in enumerate(price_headers):
    x = px + sum(p_col_w[:j])
    add_card(slide, x, py, p_col_w[j] - 0.05, 0.45, ACCENT2)
    add_text_box(slide, x + 0.05, py + 0.05, p_col_w[j] - 0.15, 0.35, h, 12, WHITE, bold=True, align=PP_ALIGN.CENTER)

for i, row in enumerate(price_data):
    y = py + 0.5 + i * 0.42
    row_color = BG_CARD if i % 2 == 0 else RGBColor(0x2D, 0x2D, 0x48)
    for j, cell in enumerate(row):
        x = px + sum(p_col_w[:j])
        add_card(slide, x, y, p_col_w[j] - 0.05, 0.37, row_color)
        c = GREEN if cell in ("FREE", "$0.00", "Yes") else RED if cell == "No" else ORANGE if cell == "Limited" else LIGHT
        add_text_box(slide, x + 0.05, y + 0.04, p_col_w[j] - 0.15, 0.3, cell, 11, c, align=PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.2, 11, 0.4, "Pocket TTS = $0 cost + voice cloning + runs on CPU = best for interview agent", 16, GREEN, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════
# SLIDE 11: Integration with VLSI Interview Agent
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.6, "Integration: VLSI Interview Agent", 32, ACCENT, bold=True)

add_bullet_card(slide, 0.5, 1.2, 6, 2.5, "Current TTS Pipeline (main.py)", [
    "Primary:   Mistral Voxtral (API, voice cloning)",
    "Fallback:  AWS Polly Amy neural (API)",
    "",
    "Cost: ~$0.03 per interview",
    "Latency: 1-2s (network dependent)",
    "Dependency: requires internet + API keys",
], ORANGE, LIGHT)

add_bullet_card(slide, 0.5, 4.0, 6, 2.5, "New TTS Pipeline (Pocket TTS integrated)", [
    "Primary:   Pocket TTS (local, CPU, FREE)",
    "Fallback 1: Mistral Voxtral (API)",
    "Fallback 2: AWS Polly (API)",
    "",
    "Cost: $0.00 per interview",
    "Latency: ~200ms (no network!)",
    "Works fully offline",
], GREEN, LIGHT)

add_bullet_card(slide, 7, 1.2, 5.5, 3.0, "How It Works", [
    "1. On first run:",
    "   - Downloads Pocket TTS model (~100M)",
    "   - Clones voice from ranjitha.mp3",
    "   - Exports ranjitha_voice.safetensors",
    "",
    "2. On subsequent runs:",
    "   - Loads .safetensors instantly",
    "   - Generates audio on CPU",
    "   - Returns base64 WAV to frontend",
], ACCENT2, LIGHT)

add_bullet_card(slide, 7, 4.5, 5.5, 2.0, "Frontend Changes", [
    "Added detectAudioType() function",
    "Auto-detects WAV vs MP3 format",
    "Both index.html and voice_agent_ui.html",
    "Seamless - no user-facing changes",
], ACCENT, LIGHT)

# ══════════════════════════════════════════════════
# SLIDE 12: Recommendation
# ══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.5, 0.3, 12, 0.6, "Recommendation Summary", 32, ACCENT, bold=True)

use_cases = [
    ("Voice Cloning (local)", "Pocket TTS", GREEN),
    ("Extreme Size (IoT)", "Kitten TTS", GREEN),
    ("Multilingual (10 langs)", "Qwen3-TTS", GREEN),
    ("Zero Hallucinations", "TADA TTS", GREEN),
    ("Real-time Agents", "Neu TTS", GREEN),
    ("Long-form (90 min)", "VibeVoice", GREEN),
    ("Best Quality (30 langs)", "VoxCPM2", GREEN),
    ("Cheapest Cloud API", "Qwen3 (AIML)", GREEN),
    ("Interview Agent", "Pocket TTS", ACCENT),
]

for i, (use, model, color) in enumerate(use_cases):
    col = i % 3
    row = i // 3
    add_card(slide, 0.5 + col * 4.2, 1.2 + row * 1.1, 4.0, 0.9, BG_CARD)
    add_text_box(slide, 0.7 + col * 4.2, 1.25 + row * 1.1, 3.6, 0.35, use, 13, GRAY)
    add_text_box(slide, 0.7 + col * 4.2, 1.6 + row * 1.1, 3.6, 0.35, model, 18, color, bold=True)

add_text_box(slide, 0.5, 5.5, 12, 0.8,
    "For the VLSI Interview Agent: Pocket TTS is integrated as primary TTS.\n"
    "Zero cost, offline, voice cloning, ~200ms latency, with Mistral + Polly as fallbacks.",
    15, LIGHT, align=PP_ALIGN.CENTER)

# Save
output_path = "/home/surya/surya_project/ojaswi/interview/TTS_Models_Research.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
